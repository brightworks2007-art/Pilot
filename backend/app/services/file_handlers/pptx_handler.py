import io

from pptx import Presentation
from pptx.util import Inches

from app.services.file_handlers.base import FileHandler


class PptxHandler(FileHandler):
    """Handles .pptx via python-pptx. Each line of content_text becomes one slide's title+body."""

    def read_text(self, file_bytes: bytes) -> str:
        prs = Presentation(io.BytesIO(file_bytes))
        lines = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    lines.append(shape.text_frame.text)
        return "\n".join(lines)

    def create(self, content_text: str) -> bytes:
        prs = Presentation()
        layout = prs.slide_layouts[1]  # Title + Content
        for line in content_text.split("\n"):
            if not line.strip():
                continue
            slide = prs.slides.add_slide(layout)
            slide.shapes.title.text = line[:80]
        buffer = io.BytesIO()
        prs.save(buffer)
        return buffer.getvalue()

    def update(self, file_bytes: bytes, instruction_result_text: str) -> bytes:
        # MVP: rebuild the deck from the new content rather than editing slides in place.
        return self.create(instruction_result_text)
